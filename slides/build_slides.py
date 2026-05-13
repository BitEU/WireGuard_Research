"""Build the CYB623 presentation deck from master_set artifacts.

Plain academic style, 16:9 widescreen. Strict layout discipline:
  * Title strip:   y = 0.30 to 1.10
  * Subtitle:      y = 1.15 to 1.50
  * Body region:   y = 1.65 to 6.65
  * Footer strip:  y = 6.85 to 7.20  (no body content past y = 6.70)

Run:
    python -m pip install python-pptx pillow
    python build_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
MASTER = os.path.join(ROOT, "master_set")
LATEX = os.path.join(ROOT, "latex-paper")
OUT = os.path.join(HERE, "WireGuard-udp2raw-detection.pptx")

CONFUSION_2F = os.path.join(MASTER, "20260509-132138_classifier_confusion.png")
ROC_2F       = os.path.join(MASTER, "20260509-132138_classifier_roc.png")
IMPORT_8F    = os.path.join(MASTER, "20260509-132129_classifier_feature_importance.png")
SCATTER      = os.path.join(LATEX,  "scatter.png")

# ---------- palette ----------
NAVY   = RGBColor(0x10, 0x2A, 0x43)
ACCENT = RGBColor(0xC9, 0x46, 0x2C)
GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF2, 0xF2, 0xF2)
BLACK  = RGBColor(0x10, 0x10, 0x10)

# ---------- layout constants (inches, 13.333 x 7.5 widescreen) ----------
SLIDE_W  = 13.333
SLIDE_H  = 7.5
LMARGIN  = 0.55
RMARGIN  = 0.55
WIDTH    = SLIDE_W - LMARGIN - RMARGIN
TITLE_TOP    = 0.30
TITLE_H      = 0.80
ACCENT_Y     = 1.10
SUBTITLE_TOP = 1.18
SUBTITLE_H   = 0.40
BODY_TOP     = 1.70
BODY_BOTTOM  = 6.55     # NOTHING past this Y in body content
FOOTER_TOP   = 6.85
FOOTER_H     = 0.30


# ---------- text primitives ----------
def _para(p, text, *, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
          space_after=4, line_spacing=None, font="Calibri"):
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(space_after)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    return run


def _add_text(slide, left, top, width, height, text, *, size=18, bold=False,
              color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    _para(tf.paragraphs[0], text, size=size, bold=bold, color=color,
          align=align, font=font)
    return box


def _bullet_list(slide, left, top, width, height, items, *, size=16,
                 color=BLACK, line_spacing=1.18, space_after=6, font="Calibri"):
    """items: list of strings or (level, string) tuples (level 0 = top)."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    first = True
    for it in items:
        if isinstance(it, tuple):
            level, txt = it
        else:
            level, txt = 0, it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        bullet = "•   " if level == 0 else "       –   "
        _para(p, bullet + txt, size=size, color=color, font=font,
              align=PP_ALIGN.LEFT, space_after=space_after,
              line_spacing=line_spacing)
    return box


def _accent_bar(slide, left, top, width, color=ACCENT, height_pt=2.5):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(height_pt))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def _title_block(slide, title, subtitle=None, *, title_size=28):
    _add_text(slide, LMARGIN, TITLE_TOP, WIDTH, TITLE_H,
              title, size=title_size, bold=True, color=NAVY,
              anchor=MSO_ANCHOR.TOP)
    _accent_bar(slide, LMARGIN, ACCENT_Y, 1.0)
    if subtitle:
        _add_text(slide, LMARGIN, SUBTITLE_TOP, WIDTH, SUBTITLE_H,
                  subtitle, size=14, color=GREY)


def _footer(slide, page_no, total):
    _add_text(slide, LMARGIN, FOOTER_TOP, 9.0, FOOTER_H,
              "Schiavone & Johannsen   ·   CYB623   ·   Pace University",
              size=10, color=GREY, align=PP_ALIGN.LEFT)
    _add_text(slide, SLIDE_W - 2.0, FOOTER_TOP, 1.5, FOOTER_H,
              f"{page_no} / {total}",
              size=10, color=GREY, align=PP_ALIGN.RIGHT)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _add_picture_fit(slide, path, left, top, max_w, max_h):
    """Add a picture, scaling to fit within (max_w x max_h) inches without
    cropping. Returns (left, top, width, height) actually used (inches)."""
    with Image.open(path) as img:
        iw, ih = img.size
    img_aspect = iw / ih
    box_aspect = max_w / max_h
    if img_aspect >= box_aspect:
        w = max_w
        h = max_w / img_aspect
    else:
        h = max_h
        w = max_h * img_aspect
    used_left = left + (max_w - w) / 2
    used_top  = top  + (max_h - h) / 2
    slide.shapes.add_picture(path,
                             Inches(used_left), Inches(used_top),
                             Inches(w), Inches(h))
    return used_left, used_top, w, h


# ---------- slide builders ----------
def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    slides = []

    # ============================================================
    # Slide 1 — Title
    # ============================================================
    s = prs.slides.add_slide(blank)
    slides.append(s)

    _add_text(s, LMARGIN, 2.20, WIDTH, 1.10,
              "Two-Feature Detection of",
              size=42, bold=True, color=NAVY)
    _add_text(s, LMARGIN, 3.00, WIDTH, 1.10,
              "WireGuard-in-udp2raw",
              size=42, bold=True, color=NAVY)
    _accent_bar(s, LMARGIN, 4.05, 2.5)
    _add_text(s, LMARGIN, 4.20, WIDTH, 0.50,
              "A Probe-and-Classify Methodology",
              size=20, color=GREY)
    _add_text(s, LMARGIN, 5.70, WIDTH, 0.40,
              "Steven Schiavone   ·   Nicholas Johannsen",
              size=18, color=BLACK)
    _add_text(s, LMARGIN, 6.10, WIDTH, 0.35,
              "Seidenberg School of CSIS · Pace University · CYB623, May 2026",
              size=12, color=GREY)
    _notes(s, """Hello, I'm Steven, this is Nicholas. We spent the semester building
a small WireGuard setup and then trying to break it.

Quick orientation before we dive in. WireGuard is a modern VPN — when
you turn it on, your traffic goes through an encrypted tunnel to a server
somewhere, and a network observer can't read what's inside. That's the
"confidentiality" part, and WireGuard does it well.

The problem we're going to talk about is different. It's not about
reading the contents. It's about whether a network observer can simply
tell that you are USING a VPN at all, even without decrypting anything.
In countries that block VPNs, that distinction is the whole game.

The story has two parts.

Part one: WireGuard is easy to spot on the wire. Every WireGuard packet
has such a recognizable shape that you can write a detector in about
twenty-five lines of Python. The Great Firewall of China actually does
this in production.

Part two: there's a popular open-source tool called udp2raw that's
supposed to disguise WireGuard. It has been downloaded roughly two
hundred and sixty thousand times. People in Iran and China rely on it.
And here's the thing — nobody has ever published a peer-reviewed security
analysis of it. Zero papers.

So we did the analysis. We came up with two ways to detect WireGuard
even when it's hiding inside udp2raw. One is "active" — we send a few
test packets at the server and watch how it replies, and we can identify
it in about eight and a half seconds. The other is "passive" — we just
watch normal traffic going by, measure two simple properties of it, and
a small machine-learning model tells WireGuard apart from regular web
traffic with ninety-nine point nine percent accuracy.

That's the talk. Let's get into it.""")

    # ============================================================
    # Slide 2 — Problem
    # ============================================================
    s = prs.slides.add_slide(blank)
    slides.append(s)
    _title_block(s,
                 "WireGuard is encrypted, but not unobservable",
                 "Strong cryptography is not the same as traffic that hides itself")

    _add_text(s, LMARGIN, BODY_TOP, WIDTH, 0.40,
              "Every WireGuard message starts with this 4-byte header:",
              size=16, color=BLACK)

    code = ("+---------+-------------------+\n"
            "| type(1) | reserved zero (3) | …rest of message\n"
            "+---------+-------------------+")
    _add_text(s, LMARGIN, 2.20, 8.5, 1.30,
              code, size=16, font="Consolas", color=NAVY)

    _bullet_list(s, LMARGIN, 3.65, WIDTH, 2.50, [
        "type byte ∈ {0x01, 0x02, 0x03, 0x04}; the next three bytes are zero.",
        "Handshake messages are exactly 148 / 92 / 64 bytes.",
        "Three concatenated checks identify a WireGuard handshake with negligible false-positive risk.",
        "A 25-line Python predicate is the full detector. We re-implement it as our baseline.",
        "Reported in production by the Great Firewall (Wu et al., USENIX Security ’23).",
    ], size=17)

    _add_text(s, LMARGIN, 5.95, WIDTH, 0.50,
              "Donenfeld’s WireGuard whitepaper (NDSS ’17) explicitly leaves obfuscation out of scope.",
              size=14, color=ACCENT, bold=True)

    _footer(s, 2, 0)
    _notes(s, """Two things on this slide. First the technical observation, then a
philosophical one that's important so we don't sound like we're attacking
the WireGuard project.

The technical observation. Every WireGuard packet starts with the exact
same four bytes. The first byte is a "message type" and it can only be
one of four values: one, two, three, or four. The next three bytes are
literally always zero. So if you're watching the network and you see a
packet whose first byte is one of those four values and the next three
are zero, that's already a really strong hint it's WireGuard.

And then it gets easier. WireGuard's "handshake" — that's the initial
back-and-forth when two computers set up the encrypted connection — uses
packets of exactly fixed sizes: a hundred and forty-eight bytes, then
ninety-two, then sixty-four. So you don't even need to decrypt anything.
You just look at the first four bytes and the size of the packet, and
you know.

You can write this whole detector in twenty-five lines of Python. We
re-implemented it ourselves as a baseline so we'd have an honest
comparison point for the rest of the paper. And this isn't theoretical —
a paper from USENIX Security 2023, by Wu and colleagues, confirms the
Great Firewall of China does exactly this in production.

Now the philosophical point, and this matters. We are NOT saying
WireGuard has a vulnerability. The author of WireGuard, Jason Donenfeld,
wrote in his original paper back in 2017 that hiding the existence of a
VPN is explicitly not WireGuard's job. His position is that if you need
to hide that you're running a VPN, you should put another tool on top of
WireGuard to handle that.

So the question we're actually asking is not "can you detect
WireGuard?" — that's a solved problem. The question is, "can you
detect WireGuard once someone has wrapped it inside the most popular
hiding tool people actually use?" That tool is udp2raw, and that's the
next slide.""")

    # ============================================================
    # Slide 3 — udp2raw and the literature gap
    # ============================================================
    s = prs.slides.add_slide(blank)
    slides.append(s)
    _title_block(s,
                 "udp2raw: the obfuscation shim, and the literature gap",
                 "Wraps WireGuard's UDP datagrams as TCP-shaped packets in userspace")

    _add_text(s, LMARGIN, BODY_TOP, WIDTH, 0.40,
              "End-to-end paths under test (same WG keys, same inner IPs; only the transport differs):",
              size=14, color=GREY)

    box_w = (WIDTH - 0.4) / 2
    # Direct path card
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(LMARGIN), Inches(2.20),
                             Inches(box_w), Inches(1.40))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY
    tf = box.text_frame
    tf.margin_left = Pt(10); tf.margin_top = Pt(8); tf.margin_right = Pt(10)
    tf.word_wrap = True
    _para(tf.paragraphs[0], "Direct (passive-detectable)",
          size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, space_after=6)
    p = tf.add_paragraph()
    _para(p, "Kali wg client", size=14, color=BLACK, font="Consolas",
          align=PP_ALIGN.CENTER, space_after=2)
    p = tf.add_paragraph()
    _para(p, "↓  UDP / 51820", size=14, color=ACCENT, font="Consolas",
          align=PP_ALIGN.CENTER, space_after=2)
    p = tf.add_paragraph()
    _para(p, "OCI wg server", size=14, color=BLACK, font="Consolas",
          align=PP_ALIGN.CENTER, space_after=2)

    # Obfuscated path card
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(LMARGIN + box_w + 0.4), Inches(2.20),
                             Inches(box_w), Inches(1.40))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = ACCENT
    tf = box.text_frame
    tf.margin_left = Pt(10); tf.margin_top = Pt(8); tf.margin_right = Pt(10)
    tf.word_wrap = True
    _para(tf.paragraphs[0], "Obfuscated (this paper's target)",
          size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, space_after=6)
    p = tf.add_paragraph()
    _para(p, "wg client  →  udp2raw client", size=13, color=BLACK,
          font="Consolas", align=PP_ALIGN.CENTER, space_after=2)
    p = tf.add_paragraph()
    _para(p, "↓  TCP / 4096  (faketcp)", size=13, color=ACCENT,
          font="Consolas", align=PP_ALIGN.CENTER, space_after=2)
    p = tf.add_paragraph()
    _para(p, "udp2raw server  →  wg server", size=13, color=BLACK,
          font="Consolas", align=PP_ALIGN.CENTER, space_after=2)

    _add_text(s, LMARGIN, 3.85, WIDTH, 0.40,
              "Why udp2raw, specifically:",
              size=17, bold=True, color=NAVY)
    _bullet_list(s, LMARGIN, 4.30, WIDTH, 2.20, [
        "≈259,000 release-binary downloads on the primary repo, "
        "+41,000 on the multiplatform fork (GitHub API, May 2026).",
        "Reported blocked by state-level DPI in Iran and elsewhere "
        "(udp2raw issue #505, net4people thread #253).",
        "Cited in the WireGuard project's own “Known Limitations” page.",
        "No peer-reviewed analysis. Searched USENIX, NDSS, IMC, FOCI, PETS, CCS.",
    ], size=15, line_spacing=1.20)

    _footer(s, 3, 0)
    _notes(s, """Two boxes at the top of the slide. They're showing the two different
ways a WireGuard client can talk to a WireGuard server in our lab. The
encryption keys are identical in both setups; the only thing that
changes is what the traffic LOOKS like to anyone watching the network.

Left box: the "direct" path. WireGuard packets just go over UDP, which
is one of the two basic ways data travels on the internet. UDP is what
WireGuard normally uses. The problem is that this is what the Great
Firewall already detects and blocks, because of the obvious fingerprint
we just talked about.

Right box: the "obfuscated" path. This is where udp2raw comes in.
udp2raw takes each WireGuard packet and re-wraps it to look like TCP
instead. TCP is the other main way data travels on the internet — it's
what your web browser uses for almost everything. The idea is that TCP
traffic is so common and so normal-looking that nobody will notice your
VPN hiding inside it. udp2raw calls this "faketcp" because it's not a
real TCP connection underneath; it's just shaped like one.

Now, why udp2raw specifically? A few reasons.

One — it's the dominant tool of its kind. Two hundred and fifty-nine
thousand downloads of the actual binary. That's not GitHub stars or
people bookmarking it, that's people downloading it to run. There's also
a fork for other platforms with another forty-one thousand downloads.

Two — there's real-world evidence that governments are blocking it. An
issue on the udp2raw GitHub from 2024 documents an Iranian ISP blocking
all of its modes within seconds. There's a long-running discussion
thread about Iran's deep packet inspection — that's the term for when a
network operator looks at the contents of packets, not just where
they're going — and udp2raw is one of the tools that gets blocked.

Three, and this is the gap we're filling — nobody has academically
analyzed it. We searched the major security conferences. USENIX, NDSS,
IMC, FOCI, PETS, CCS. These are the venues where this kind of work
gets published. Nothing on udp2raw. So this is genuinely an unstudied
target.

One last thing — you might wonder why we don't also test other tools
like AmneziaWG or swgp-go. The answer is that those tools modify
WireGuard itself; they're a different category of fix. udp2raw is the
only big, popular, unstudied tool in the "wrap it generically" category,
so we focused there.""")

    # ============================================================
    # Slide 4 — Methodology overview
    # ============================================================
    s = prs.slides.add_slide(blank)
    slides.append(s)
    _title_block(s,
                 "Two attacks, one lab",
                 "Following Xue et al.'s active-then-passive pattern (USENIX Security ’22) for OpenVPN")

    box_w = (WIDTH - 0.5) / 2
    # C1 card
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(LMARGIN), Inches(BODY_TOP),
                             Inches(box_w), Inches(4.50))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY
    tf = box.text_frame
    tf.margin_left = Pt(16); tf.margin_top = Pt(14); tf.margin_right = Pt(16)
    tf.word_wrap = True
    _para(tf.paragraphs[0], "C1 · Active prober",
          size=22, bold=True, color=NAVY, space_after=10)
    for line in [
        "Send 5 stateless TCP probes to (IP, port).",
        "Real Linux TCP answers all 5.",
        "udp2raw answers SYN, drops 3 of 4 follow-ups.",
        "Verdict in ≤ 9 seconds. Same verdict 5/5 trials.",
    ]:
        p = tf.add_paragraph()
        _para(p, "•   " + line, size=14, color=BLACK, space_after=6,
              line_spacing=1.18)
    p = tf.add_paragraph(); p.space_before = Pt(12)
    _para(p, "Adversary cost:  ≈ 1.7 KB and ≈ 9 s per host.",
          size=12, color=GREY, space_after=2)

    # C2 card
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(LMARGIN + box_w + 0.5), Inches(BODY_TOP),
                             Inches(box_w), Inches(4.50))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = ACCENT
    tf = box.text_frame
    tf.margin_left = Pt(16); tf.margin_top = Pt(14); tf.margin_right = Pt(16)
    tf.word_wrap = True
    _para(tf.paragraphs[0], "C2 · Two-feature classifier",
          size=22, bold=True, color=ACCENT, space_after=10)
    for line in [
        "bulk_fraction:  share of MTU-class packets.",
        "ack60_fraction:  share of 60-byte TCP ACKs.",
        "Both fall out of how the protocols are written, "
        "not statistical luck.",
        "Random Forest:  99.9 % accuracy, AUC = 1.000.",
    ]:
        p = tf.add_paragraph()
        _para(p, "•   " + line, size=14, color=BLACK, space_after=6,
              line_spacing=1.18)
    p = tf.add_paragraph(); p.space_before = Pt(12)
    _para(p, "Trained on 24 WG flows + 1058 ISCXVPN2016 non-VPN flows.",
          size=12, color=GREY, space_after=2)

    _footer(s, 4, 0)
    _notes(s, """Quick framing on this slide. There are two general ways to detect a
hidden service on a network. "Active" means you send something at the
server and watch how it responds. "Passive" means you don't send
anything — you just watch traffic going by and look for patterns. We do
both.

We didn't invent this two-pronged approach. A team at the University of
Michigan led by Diwen Xue published a paper at USENIX Security 2022 that
won "Distinguished Paper" — that's the conference's top award — where
they did exactly this for a different VPN called OpenVPN. We're applying
their same playbook to WireGuard-plus-udp2raw, which their paper didn't
cover.

Left box, contribution one: the active prober. We send five carefully
chosen TCP packets at the target server. Five packets total. The whole
test takes under nine seconds, and we always get the same answer when
we re-run it. I'll walk through the actual packets on the next slide.

Right box, contribution two: the passive classifier. We just watch
traffic and measure two things about each "flow." A flow is a
conversation between two computers — one IP address talking to another
on specific ports. So for every flow we see, we compute these two
numbers, and a small machine learning model uses them to decide:
WireGuard, or not?

The two numbers are called "bulk_fraction" and "ack60_fraction." Don't
worry about what they mean yet — slide six is dedicated to explaining
them. The key thing is they're not random statistical patterns we got
lucky with. They come directly from how WireGuard and udp2raw are
written. WireGuard always pads its packets to the maximum size, and
udp2raw sends back one acknowledgment packet for every data packet it
receives. Those two implementation choices are essentially impossible to
hide.

The model we use is called a Random Forest. Think of it as a system
that learns rules from examples — you give it a bunch of WireGuard
flows and a bunch of normal flows, and it figures out the boundary
between them. On our data, it's correct ninety-nine point nine percent
of the time.

For training data, we captured twenty-four WireGuard flows ourselves in
the lab. For the "normal traffic" comparison, we used a public research
dataset called ISCXVPN2016, which has over a thousand recordings of
ordinary web browsing, video, email, chat, and so on. It's the standard
dataset everyone uses for this kind of work.""")

    # ============================================================
    # Slide 5 — C1 results
    # ============================================================
    s = prs.slides.add_slide(blank)
    slides.append(s)
    _title_block(s,
                 "C1: one set of probes distinguishes udp2raw from real TCP",
                 "Three targets on the same OCI host. Same verdict in 5/5 repeated trials.")

    headers = ["Probe", "TCP/22  (real SSH)", "TCP/4096  (udp2raw)", "TCP/9999  (closed)"]
    rows = [
        ["P1   SYN",                "SYN-ACK in 181 ms",   "SYN-ACK in 163 ms",   "no response"],
        ["P2   SYN + 32 B payload", "SYN-ACK + 2 B echo",  "silent",              "no response"],
        ["P3   bogus-seq ACK",      "RST in 81 ms",        "RST in 91 ms",        "no RST"],
        ["P4   SYN, WS = 14",       "SYN-ACK",             "silent",              "no response"],
        ["P5   FIN to unopened",    "RST",                 "silent",              "no response"],
    ]

    top = BODY_TOP
    col_lefts  = [LMARGIN, 3.40, 6.50, 9.80]
    col_widths = [2.85,    3.10, 3.30, 2.85]
    row_h = 0.42

    for i, h in enumerate(headers):
        _add_text(s, col_lefts[i], top, col_widths[i], row_h,
                  h, size=13, bold=True, color=NAVY)
    _accent_bar(s, LMARGIN, top + row_h, WIDTH, color=NAVY)

    for r, row in enumerate(rows):
        y = top + row_h + 0.10 + r * 0.40
        for i, val in enumerate(row):
            color = ACCENT if val == "silent" else BLACK
            font = "Consolas" if i > 0 else "Calibri"
            _add_text(s, col_lefts[i], y, col_widths[i], row_h,
                      val, size=13, color=color, font=font)

    # Summary band
    y = top + row_h + 0.10 + len(rows) * 0.40 + 0.15
    _accent_bar(s, LMARGIN, y, WIDTH, color=NAVY)
    sum_y = y + 0.18

    _add_text(s, col_lefts[0], sum_y, col_widths[0], row_h,
              "Total time", size=13, bold=True, color=NAVY)
    for i, val in enumerate(["0.5 s", "8.5 s", "15.5 s"]):
        _add_text(s, col_lefts[i+1], sum_y, col_widths[i+1], row_h,
                  val, size=13, bold=True, color=BLACK, font="Consolas")
    sum_y2 = sum_y + 0.40
    _add_text(s, col_lefts[0], sum_y2, col_widths[0], row_h,
              "Verdict", size=13, bold=True, color=NAVY)
    for i, val in enumerate(["REAL_TCP", "UDP2RAW_FAKETCP", "UNREACHABLE"]):
        c = ACCENT if "UDP2RAW" in val else BLACK
        _add_text(s, col_lefts[i+1], sum_y2, col_widths[i+1], row_h,
                  val, size=13, bold=True, color=c, font="Consolas")

    _add_text(s, LMARGIN, 6.20, WIDTH, 0.40,
              "Discriminator: count silent responses among {P2, P4, P5}.  Real TCP answers all three.  udp2raw answers none.",
              size=14, color=GREY)

    _footer(s, 5, 0)
    _notes(s, """This is the active prober's results. Three columns — same server, three
different ports. A "port" is just a number that identifies which
service on a server you're talking to. Port twenty-two is real SSH,
which is normal remote-login software running on the kernel — that's
our "real TCP" control. Port four thousand ninety-six is where udp2raw
is listening. Port nine thousand nine hundred ninety-nine is closed,
nothing running there — that's our negative control.

Quick TCP refresher so the table reads cleanly. When two computers want
to start a TCP connection, the client sends a "SYN" packet — that's
short for synchronize, it's basically "hello, want to talk?" The server
replies with a "SYN-ACK" — "hello back, sure." That's the start of every
TCP connection. There's also "RST" — short for reset — which is what a
server sends to say "no, go away." And "FIN" — for finish — which means
"I'm done talking." Those four — SYN, SYN-ACK, RST, FIN — are most of
what you need to read this table.

Five rows in the table, one per probe. Look at the orange "silent"
entries in the middle column.

The first probe is just a normal SYN. udp2raw answers it. Of course it
does — that's the whole point of the disguise. If it didn't answer SYNs
it wouldn't look like a TCP server at all.

The third probe is a bogus ACK — we pretend to acknowledge a connection
that doesn't exist. udp2raw also answers this one, but here's the
subtle reason: udp2raw itself doesn't generate the response, the
underlying Linux kernel does, because udp2raw didn't recognize the
packet and the kernel emits its own reset by default. So probe one and
probe three don't help us discriminate.

The interesting ones are probes two, four, and five. Probe two is a SYN
packet with thirty-two bytes of data attached. That's unusual but legal,
and a real Linux kernel always replies. Probe four uses an unusual TCP
option — specifically a "window scale" of fourteen, which is at the edge
of what's allowed. A real kernel still replies. Probe five is a FIN
packet — a "goodbye" packet — sent to a port where no conversation was
ever started. Real kernels reply to this too.

udp2raw replies to none of them. Silent, silent, silent. That's the
giveaway. udp2raw was written to handle one specific TCP pattern, and
anything outside that pattern, it just drops on the floor. A real Linux
kernel never drops these — it always reacts.

So our rule is simple: count how many of probes two, four, and five
were silent. If it's three out of three, it's udp2raw. Real TCP scores
zero out of three. Done.

The bottom of the table is the timing. Half a second for real TCP. Eight
and a half seconds for udp2raw. The reason udp2raw takes longer is we
wait the full timeout on each of the silent probes. The closed port
takes fifteen and a half because there's nothing replying to anything.
And we ran this against udp2raw five separate times — same verdict every
time. Five out of five.""")

    # ============================================================
    # Slide 6 — C2 features (table + scatter)
    # ============================================================
    s = prs.slides.add_slide(blank)
    slides.append(s)
    _title_block(s,
                 "C2: two features carry the entire signal",
                 "Class means across 24 WireGuard flows and 1058 ISCXVPN2016 non-VPN flows")

    # Left half: feature table
    headers = ["Feature", "WireGuard", "Non-VPN", "Cohen's d"]
    rows_top = [
        ["bulk_fraction",   "0.553 ± 0.186", "0.019 ± 0.066", "3.83"],
        ["ack60_fraction",  "0.295 ± 0.187", "0.069 ± 0.169", "1.27"],
    ]
    rows_dim = [
        ["len_entropy",          "1.94 ± 0.82", "1.23 ± 1.42",  "—"],
        ["dominant_size_frac.",  "0.50 ± 0.19", "0.71 ± 0.30",  "(inv.)"],
        ["top3_size_fraction",   "0.88 ± 0.15", "0.87 ± 0.19",  "—"],
        ["rate_pps",             "3.23 ± 0.90", "3.37 ± 30.27", "—"],
    ]

    top = BODY_TOP
    col_lefts  = [LMARGIN, 2.70, 4.70, 6.40]
    col_widths = [2.10,    1.95, 1.65, 0.85]

    for i, h in enumerate(headers):
        _add_text(s, col_lefts[i], top, col_widths[i], 0.32,
                  h, size=12, bold=True, color=NAVY)
    _accent_bar(s, LMARGIN, top + 0.32, 6.85, color=NAVY)

    y = top + 0.42
    for row in rows_top:
        for i, val in enumerate(row):
            font = "Consolas" if i > 0 else "Calibri"
            color = ACCENT if i == 3 else BLACK
            bold = (i == 3)
            _add_text(s, col_lefts[i], y, col_widths[i], 0.30,
                      val, size=12, color=color, bold=bold, font=font)
        y += 0.36

    _accent_bar(s, LMARGIN, y + 0.04, 6.85, color=GREY, height_pt=1.0)
    y += 0.16

    for row in rows_dim:
        for i, val in enumerate(row):
            font = "Consolas" if i > 0 else "Calibri"
            _add_text(s, col_lefts[i], y, col_widths[i], 0.30,
                      val, size=11, color=GREY, font=font)
        y += 0.32

    # Right half: scatter plot, fitted to a strict box
    SCATTER_BOX_LEFT = 7.55
    SCATTER_BOX_TOP  = BODY_TOP - 0.05
    SCATTER_BOX_W    = SLIDE_W - SCATTER_BOX_LEFT - RMARGIN
    SCATTER_BOX_H    = 4.20
    if os.path.exists(SCATTER):
        _add_picture_fit(s, SCATTER,
                         SCATTER_BOX_LEFT, SCATTER_BOX_TOP,
                         SCATTER_BOX_W, SCATTER_BOX_H)
    _add_text(s, SCATTER_BOX_LEFT, SCATTER_BOX_TOP + SCATTER_BOX_H + 0.05,
              SCATTER_BOX_W, 0.50,
              "All 1082 flows in the two-feature space. "
              "24/24 WG flows have bulk_fraction ≥ 0.32; only 4/1058 "
              "non-VPN flows fall in the (≥0.2, ≥0.15) corner.",
              size=10, color=GREY)

    _add_text(s, LMARGIN, 6.05, 6.85, 0.50,
              "Naive intuition picked dominant_size_fraction. "
              "Real-world background INVERTS that signal "
              "(many idle flows are degenerately concentrated).",
              size=11, color=ACCENT)

    _footer(s, 6, 0)
    _notes(s, """OK this is the heart of the passive classifier. Two numbers we measure
about every flow.

First feature, called "bulk_fraction." This is just the fraction of
packets in the flow that are "big." We defined "big" as twelve hundred
bytes or more. The reason is something called MTU — short for Maximum
Transmission Unit, which is the largest packet size a network link will
carry, usually around fifteen hundred bytes. WireGuard pads every single
one of its packets up to right near that maximum. Web browsers don't do
this. A typical web flow has a mix of big packets when you're
downloading something and tiny packets the rest of the time. So:

WireGuard flows: about fifty-five percent of their packets are big.
Normal background traffic: about two percent. That's a huge gap.

The table also lists something called "Cohen's d." This is a statistics
term that measures how cleanly separated two groups of numbers are. It's
on a scale where zero point eight is already considered a "large"
difference. For bulk_fraction we get three point eight three — almost
five times the threshold for "large." This is a giant, obvious
difference, not a subtle one we squeezed out with statistics.

Second feature, "ack60_fraction." This one's a little more intricate so
bear with me. When you make a TCP connection, every time you send data,
the other side sends back a tiny acknowledgment packet — sixty bytes,
usually. These are called ACKs. In normal bulk file transfers, the
receiving side waits and sends one ACK for every TWO data packets, to
save bandwidth. That's called "delayed ACK." So normal TCP has about
thirty-three percent ACKs.

But udp2raw breaks this. udp2raw takes each WireGuard packet and wraps
it as one separate TCP segment. The receiving computer's kernel sees
these as separate arrivals and ACKs each one individually — no delayed
ACK happens. So you end up with one ACK per data packet, fifty-fifty,
and the ACK fraction climbs.

For WireGuard-in-udp2raw, ack60_fraction is about thirty percent. For
ordinary traffic, it's about seven percent. Cohen's d of one point two
seven — still a large effect.

The scatter plot on the right is showing all one thousand eighty-two
flows we have, plotted on these two axes. Orange dots are WireGuard.
Grey dots are normal traffic. You can see they live in completely
different regions of the plot. They barely overlap. All twenty-four of
our WireGuard flows have a bulk_fraction of at least zero point three
two. Out of one thousand fifty-eight normal flows, only four are
anywhere near the WireGuard cluster.

Bottom-left in orange, a quick aside about scientific honesty. Going in,
I expected a different feature called "dominant_size_fraction" — meaning
how often the most-common packet size shows up — to be the smoking gun.
Turned out the data went the opposite direction from what I expected.
The reason is that the public dataset contains lots of idle "keepalive"
flows where every packet is literally identical, so they score higher on
that feature than WireGuard does. The real-world data corrected my bad
guess. We dropped that feature and three others, and the two that
survived are the ones grounded in actual protocol mechanics, not in
intuition.""")

    # ============================================================
    # Slide 7 — C2 classifier results
    # ============================================================
    s = prs.slides.add_slide(blank)
    slides.append(s)
    _title_block(s,
                 "C2: two-feature Random Forest hits 99.9 %",
                 "5-fold stratified cross-validation; class-balanced; n = 1082")

    # Left: the four big numbers
    big_y = BODY_TOP + 0.10
    _add_text(s, LMARGIN, big_y, 5.5, 0.40,
              "Two-feature Random Forest", size=18, bold=True, color=NAVY)

    metrics = [
        ("99.91 %", "Accuracy"),
        ("1.000",   "Precision  (zero false positives in 1058)"),
        ("0.958",   "Recall  (23 of 24 WG flows caught)"),
        ("1.000",   "ROC AUC"),
    ]
    y = big_y + 0.55
    for big, label in metrics:
        _add_text(s, LMARGIN, y, 2.30, 0.55,
                  big, size=28, bold=True, color=ACCENT, font="Consolas")
        _add_text(s, LMARGIN + 2.45, y + 0.12, 3.60, 0.45,
                  label, size=13, color=BLACK)
        y += 0.62

    # Comparison table
    cmp_y = y + 0.20
    _accent_bar(s, LMARGIN, cmp_y, 6.0, color=NAVY)
    _add_text(s, LMARGIN, cmp_y + 0.10, 6.0, 0.32,
              "Adding more features does not help.",
              size=13, bold=True, color=NAVY)
    rows = [
        ["RF, 2 features", "0.999", "1.000", "0.958"],
        ["RF, 8 features", "0.998", "1.000", "0.917"],
        ["LR, 2 features", "0.969", "0.414", "1.000"],
    ]
    headers = ["model", "acc", "prec", "rec"]
    cl = [LMARGIN, 3.10, 4.10, 5.10]
    cw = [2.50, 1.00, 1.00, 1.00]
    yy = cmp_y + 0.50
    for i, h in enumerate(headers):
        _add_text(s, cl[i], yy, cw[i], 0.30,
                  h, size=10, bold=True, color=GREY)
    yy += 0.30
    for row in rows:
        for i, v in enumerate(row):
            font = "Consolas" if i > 0 else "Calibri"
            _add_text(s, cl[i], yy, cw[i], 0.30,
                      v, size=11, color=BLACK, font=font)
        yy += 0.30

    # Right: confusion + ROC, fitted into strict boxes
    PLOT_LEFT = 7.55
    PLOT_W    = SLIDE_W - PLOT_LEFT - RMARGIN  # ~5.2 inches
    PLOT_H    = (BODY_BOTTOM - BODY_TOP - 0.20) / 2  # ~2.3 inches each
    if os.path.exists(CONFUSION_2F):
        _add_picture_fit(s, CONFUSION_2F,
                         PLOT_LEFT, BODY_TOP,
                         PLOT_W, PLOT_H)
    if os.path.exists(ROC_2F):
        _add_picture_fit(s, ROC_2F,
                         PLOT_LEFT, BODY_TOP + PLOT_H + 0.20,
                         PLOT_W, PLOT_H)

    _footer(s, 7, 0)
    _notes(s, """OK so we took those two features from the last slide and trained the
classifier. Here are the results.

Before reading the numbers, let me quickly define them, because each one
measures something different.

"Accuracy" is the simplest — out of every flow we tested, how many did
we label correctly? We got ninety-nine point nine one percent. So out of
about a thousand flows, we mislabeled around one.

"Precision" answers the question: when the classifier says "this IS
WireGuard," how often is it right? We got one point zero. That's perfect.
Zero false alarms. We never accused a normal web flow of being WireGuard.

"Recall" is the opposite question: of all the real WireGuard flows that
existed, how many did we catch? Zero point nine five eight. That means
we caught twenty-three of the twenty-four WireGuard flows. One slipped
through. So we're slightly conservative — we miss the occasional one,
but we never falsely accuse.

"ROC AUC" is the trickiest one to explain. Imagine you could tune the
classifier to be more paranoid — willing to call more flows WireGuard —
or more relaxed. Each setting trades false alarms for missed WireGuards.
ROC AUC is a single number that summarizes how well the classifier does
across ALL of those tunings at once. The scale is zero point five for
random guessing, one point zero for a perfect classifier. We got
one point zero. Couldn't ask for better.

Now the smaller table below — this is showing what happens if we add
more features. The "RF, two features" row is the headline classifier we
just discussed. The "RF, eight features" row added six more measurements:
things like packet-size variety, how concentrated the sizes are, how
fast packets are flowing, and so on. You'd think more information means
better results. It doesn't. Recall actually drops a tiny bit — we catch
twenty-two flows instead of twenty-three. The extra features added noise
without adding signal. The data was telling us "those two features were
enough."

"LR" in the third row stands for Logistic Regression, which is a simpler
type of model than Random Forest — it just tries to draw a straight line
between the two groups. It catches every single WireGuard flow but it
also raises about three percent false alarms, because the boundary
between the two groups isn't actually a straight line.

The pictures on the right. The top one is a "confusion matrix" — it's a
two-by-two grid showing how many WireGuard flows got correctly labeled
as WireGuard, how many normal flows got correctly labeled as normal,
and the mistakes in the other two corners. Almost everything is on the
diagonal, which means almost everything is correct.

The bottom picture is the "ROC curve" — that's where the ROC AUC number
comes from. The curve hugs the top-left corner of the plot, which is
what a near-perfect classifier looks like. A diagonal line through the
middle would mean random guessing. We're nowhere near that.""")

    # ============================================================
    # Slide 8 — Related work and limitations
    # ============================================================
    s = prs.slides.add_slide(blank)
    slides.append(s)
    _title_block(s,
                 "Where this sits in the literature",
                 "Narrow contribution alongside the recent Xue et al. line of work")

    col_w = (WIDTH - 0.5) / 2
    _add_text(s, LMARGIN, BODY_TOP, col_w, 0.38,
              "Closely related",
              size=17, bold=True, color=NAVY)
    _bullet_list(s, LMARGIN, BODY_TOP + 0.45, col_w, 4.20, [
        "Xue et al., USENIX Sec ’22.  Passive + active for OpenVPN. Methodological template.",
        "Xue et al., USENIX Sec ’24.  Encapsulated TLS handshakes. Doesn’t fire on WireGuard (no nested TLS).",
        "Xue et al., NDSS ’25.  Cross-layer RTT. Attenuated on UDP-only short flows.",
        "Wu et al., USENIX Sec ’23.  GFW entropy detector for fully-encrypted protocols.",
        "Alice et al., IMC ’20.  Shadowsocks active probes. Probe taxonomy.",
        "Frolov & Wustrow, NDSS ’20.  Probe-resistant proxy detection.",
    ], size=13, line_spacing=1.16, space_after=6)

    right_x = LMARGIN + col_w + 0.5
    _add_text(s, right_x, BODY_TOP, col_w, 0.38,
              "What this paper does NOT claim",
              size=17, bold=True, color=ACCENT)
    _bullet_list(s, right_x, BODY_TOP + 0.45, col_w, 4.20, [
        "Single client/server pair. No ISP-scale deployment trace.",
        "udp2raw v20230206.0 only. Later forks may behave differently.",
        "Static defenses: a delayed-ACK patch + WG MTU below the bulk threshold weakens C2.",
        "C1’s P3 (bogus ACK) is mildly hostile. Adversary should fire only at already-suspected hosts.",
        "Not benchmarked against AmneziaWG, swgp-go, wstunnel. Future work.",
    ], size=13, line_spacing=1.16, space_after=6)

    _footer(s, 8, 0)
    _notes(s, """This slide does two things. The left side puts us in conversation with
the existing research, so you know we didn't reinvent the wheel. The
right side is us being honest about what we did NOT do.

Left side. The dominant research group in this area is at the University
of Michigan, led by Diwen Xue. They have three relevant papers.

Their 2022 paper at USENIX Security broke obfuscated OpenVPN. That's
where we got our two-stage active-plus-passive methodology. OpenVPN is
a different VPN than WireGuard, and the specific attacks don't carry
over, but the playbook does.

Their 2024 paper looks for nested encryption handshakes — basically,
when a tool tries to disguise traffic as HTTPS, there's often a real
TLS handshake hidden inside. They look for that. It doesn't work
against WireGuard, because WireGuard doesn't use TLS at all. There's
no nested handshake to find.

Their 2025 paper at NDSS uses timing differences between layers of the
network stack — basically measuring round-trip times at different
levels and looking for inconsistencies. It works well on long,
high-bandwidth connections. Our captures are short and use UDP, where
that timing signal is weak. So that approach doesn't really apply to
our case either.

The bottom line is: those general-purpose techniques exist, but none of
them solve the WireGuard-in-udp2raw case specifically. So the
protocol-specific approach we took — features that come from how
WireGuard and udp2raw are actually written — is the available
approach.

We also cite Wu et al. on the Great Firewall's entropy filter, and
older work on Shadowsocks and TLS-based circumvention tools. That gives
us the broader context.

Right side. What we are NOT claiming.

One — we ran this in a lab, with one client and one server. We don't
have data from actual internet service providers at scale. So we can't
say "this would work on every network in the world." It's possible some
of our signals weaken in environments we didn't test.

Two — we tested one specific version of udp2raw, the February 2023
release. There are forks of udp2raw with slightly different behavior
that we didn't test.

Three — both of our attacks could be defeated by patching udp2raw. If
the maintainer added delayed-ACK behavior and recommended a smaller
WireGuard packet size, our passive classifier would lose most of its
edge. If they made the TCP emulation more thorough, our active prober
would lose its edge too. We're describing the current state of the
tool, not making a permanent claim.

Four — the active prober's third probe, the bogus ACK, is mildly
hostile. In a real deployment you'd only fire that one at hosts you
already suspect.

Five — we didn't compare against the other obfuscation tools like
AmneziaWG, swgp-go, or wstunnel. That's left for future work.

The point of this slide is to be honest about scope so people can trust
the parts we DID prove.""")

    # ============================================================
    # Slide 9 — Takeaways
    # ============================================================
    s = prs.slides.add_slide(blank)
    slides.append(s)
    _title_block(s, "Takeaway")

    msgs = [
        ("WireGuard's wire format leaks.", NAVY,
         "A 25-line classifier handles bare WG. udp2raw is the deployed fix."),
        ("udp2raw doesn't actually fix it.", ACCENT,
         "Five TCP probes in 8.5 seconds (C1).  Two flow features, RF, AUC = 1.0 (C2)."),
        ("Both signals are structural.", NAVY,
         "udp2raw drops unfamiliar TCP shapes; WG pads to MTU; udp2raw acks 1:1. None are statistical luck."),
        ("Defenses are obvious; the upstream tool does not implement them.", ACCENT,
         "Delayed-ACK plus a sub-bulk WG MTU closes C2. More-conformant TCP emulation closes C1."),
    ]

    y = BODY_TOP
    for headline, color, body in msgs:
        _add_text(s, LMARGIN, y, WIDTH, 0.40,
                  headline, size=20, bold=True, color=color)
        _add_text(s, LMARGIN, y + 0.45, WIDTH, 0.50,
                  body, size=13, color=BLACK)
        y += 1.10

    _add_text(s, LMARGIN, 6.30, WIDTH, 0.30,
              "Code, configs, evidence, paper, and slides are in the project repository.",
              size=12, color=GREY)

    _footer(s, 9, 0)
    _notes(s, """Four takeaways. If you remember nothing else from this talk, remember
these.

Number one. WireGuard's packets are easy to recognize on the network.
That's not news — researchers and the Great Firewall have known this
for years. We re-built the detector ourselves, in twenty-five lines of
Python, just so we have a fair baseline to compare against.

Number two. The most popular tool people use to hide WireGuard — udp2raw,
which has been downloaded a quarter million times — doesn't actually
hide it well. You can spot it actively in eight and a half seconds by
sending five test packets. Or you can spot it passively, just watching
traffic, with a tiny machine-learning model that's right essentially a
hundred percent of the time. Take your pick.

Number three. These results are not statistical flukes or lucky
patterns we found in one dataset. Both detectors come from how the
software is fundamentally written. WireGuard always pads its packets to
the maximum size — that's a deliberate design choice in the protocol.
udp2raw always acknowledges every wrapped packet one-for-one, because
that's how its wrapping mechanism works. And udp2raw silently drops
weird TCP packets because the person who wrote it only thought about
the one specific TCP pattern. These are mechanical consequences of the
code, not statistical accidents — so they should transfer to other
networks and other situations.

Number four. The defenses are not hard. If the udp2raw author added
delayed acknowledgments, and recommended users set a smaller WireGuard
packet size, the passive classifier mostly stops working. If they made
the TCP emulation more thorough — actually responding to weird probes
the way a real Linux kernel does — the active prober stops working.
We describe both fixes in the paper. We didn't actually write the patch
ourselves and submit it upstream. That's a pull request someone else
can write.

That's the talk. Happy to take questions, and we have a live demo
ready if there's interest.""")

    # ============================================================
    # Patch footers with actual total
    # ============================================================
    total = len(slides)
    for i, s in enumerate(slides, 1):
        # remove placeholder footers (the "X / 0" boxes added by _footer)
        for sh in list(s.shapes):
            if sh.has_text_frame:
                txt = sh.text_frame.text
                if " / 0" in txt:
                    sh._element.getparent().remove(sh._element)
        _footer(s, i, total)

    prs.save(OUT)
    print(f"wrote {OUT}  ({total} slides)")


if __name__ == "__main__":
    build()
